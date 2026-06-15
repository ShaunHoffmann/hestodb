from pathlib import Path
import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from hestodb.util import find_latest_report_pptx, parse_file_path


def _create_report_file(
    root: Path, year: str, project: str, filename: str, mtime: float, date: str = "4/30/2026"
) -> Path:
    """Create a mock PowerPoint file with a summary slide containing a date field."""
    report_dir = root / year / project / "Reports" / "HESTO"
    report_dir.mkdir(parents=True, exist_ok=True)
    file_path = report_dir / filename
    
    # Create a presentation with a summary slide
    prs = Presentation()
    # Use a layout with title and content (index 1) instead of blank
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set the title using the title placeholder
    title = slide.shapes.title
    title.text = "Summary"
    
    # Add a table with summary data
    rows, cols = 8, 2
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(4)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Add headers and sample data
    summary_data = [
        ("Date", date),
        ("Proposal ID", "24-HTIDS24-0009"),
        ("Principal Investigator", "Jane Smith"),
        ("Affiliation", "Test University"),
        ("Research Regime", "Magnetosphere"),
        ("Technology Name", "TestTech"),
        ("Project Summary", "Test summary"),
        ("Proposal Title", "Test Proposal"),
    ]
    
    for i, (key, value) in enumerate(summary_data):
        table_shape.cell(i, 0).text = key
        table_shape.cell(i, 1).text = value
    
    prs.save(str(file_path))
    
    # Make mtimes deterministic for newest-file selection tests.
    import os
    os.utime(file_path, (mtime, mtime))
    return file_path


def test_parse_file_path_extracts_expected_metadata_from_string():
    file_path = "root/24/24-HTIDS24-0009 Jane Smith/Reports/HESTO/report.pptx"
    metadata = parse_file_path(file_path)

    assert metadata["year"] == 2024
    assert metadata["project_id"] == "24-HTIDS24-0009"
    assert metadata["principal_investigator"] == "Jane Smith"
    assert metadata["filename"] == "report.pptx"


def test_parse_file_path_extracts_expected_metadata_from_path():
    file_path = Path("root/25/25-HTIDS24-0010 Ada Lovelace/Reports/HESTO/update.pptx")
    metadata = parse_file_path(file_path)

    assert metadata["year"] == 2025
    assert metadata["project_id"] == "25-HTIDS24-0010"
    assert metadata["principal_investigator"] == "Ada Lovelace"
    assert metadata["filename"] == "update.pptx"


def test_find_latest_report_pptx_returns_empty_dataframe_when_no_matches(tmp_path):
    result = find_latest_report_pptx(tmp_path)

    assert result.empty
    assert result.index.name == "filename"
    assert list(result.columns) == [
        "file_path",
        "project_id",
        "report_date",
        "principal_investigator",
        "year",
        "modified",
        "folder",
        "is_active",
    ]


def test_find_latest_report_pptx_selects_newest_file_per_folder_and_skips_lock_files(
    tmp_path,
):
    project = "24-HTIDS24-0009 Jane Smith"

    older = _create_report_file(
        tmp_path,
        "24",
        project,
        "HESTO-202606_Report_24-HTIDS24-0009_old.pptx",
        mtime=1700000000,
        date="6/1/2026",
    )
    newer = _create_report_file(
        tmp_path,
        "24",
        project,
        "HESTO-202607_Report_24-HTIDS24-0009_new.pptx",
        mtime=1700000100,
        date="7/1/2026",
    )
    _create_report_file(
        tmp_path,
        "24",
        project,
        "~$HESTO-202402_Report_24-HTIDS24-0009_new.pptx",
        mtime=1700000200,
    )

    result = find_latest_report_pptx(str(tmp_path))

    assert list(result.index) == [newer.name]
    assert result.loc[newer.name, "file_path"] == newer
    assert result.loc[newer.name, "project_id"] == "24-HTIDS24-0009"
    assert result.loc[newer.name, "principal_investigator"] == "Jane Smith"
    assert result.loc[newer.name, "year"] == 2024
    assert result.loc[newer.name, "folder"] == newer.parent
    assert older.name not in result.index


def test_find_latest_report_pptx_sorts_rows_by_year_descending(tmp_path):
    recent = _create_report_file(
        tmp_path,
        "25",
        "25-HTIDS24-0010 Ada Lovelace",
        "HESTO-202607_a_report.pptx",
        mtime=1700000300,
        date="7/15/2026",
    )
    older = _create_report_file(
        tmp_path,
        "24",
        "24-HTIDS24-0009 Jane Smith",
        "HESTO-202606_b_report.pptx",
        mtime=1700000400,
        date="6/15/2026",
    )

    result = find_latest_report_pptx(tmp_path)

    assert list(result.index) == [recent.name, older.name]
    assert list(result["year"]) == [2025, 2024]


def test_find_latest_report_pptx_mocks_notebook_htides_directory_layout(tmp_path):
    # Mirrors the notebook call:
    # find_latest_report_pptx(Path(".../Box/HESTO/HTIDeS"))
    notebook_root = tmp_path / "Box" / "HESTO" / "HTIDeS"

    _create_report_file(
        notebook_root,
        "24",
        "24-HTIDS24-0009 Jane Smith",
        "HESTO-202612_Report_24-HTIDS24_Smith-rev3.pptx",
        mtime=1700000100,
        date="12/15/2026",
    )
    _create_report_file(
        notebook_root,
        "24",
        "24-HTIDS24-0010 Ada Lovelace",
        "HESTO-202612_Report_24-HTIDS24_Lovelace-rev1.pptx",
        mtime=1700000200,
        date="12/20/2026",
    )
    _create_report_file(
        notebook_root,
        "25",
        "25-HTIDS25-0001 Grace Hopper",
        "HESTO-202701_Report_25-HTIDS25_Hopper-rev1.pptx",
        mtime=1700000300,
        date="1/15/2027",
    )

    files = find_latest_report_pptx(notebook_root)

    assert len(files) == 3
    assert len(files[files["year"] == 2024]) == 2
    assert files.iloc[0]["year"] == 2025
    assert all(path.name == "HESTO" for path in files["folder"])


@pytest.mark.parametrize("bad_path", ["report.pptx", "root/24/report.pptx"])
def test_parse_file_path_raises_for_unexpected_structure(bad_path):
    with pytest.raises((ValueError, IndexError)):
        parse_file_path(bad_path)
