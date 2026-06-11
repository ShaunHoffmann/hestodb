from dataclasses import dataclass, field
from pathlib import Path
from pptx import Presentation
from datetime import datetime
import csv
import re
import logging
import math
import pandas as pd

logger = logging.getLogger(__name__)

RESEARCH_REGIMES = [
    "magnetosphere",
    "itm",
    "solar science",
    "space weather",
    "heliosphere",
]
STATUS_VALUES = ["g", "y", "r", "green", "yellow", "red"]
PUBLICATION_TYPES = [
    "Peer-reviewed publication",
    "Non-Peer-reviewed publication",
    "Conference presentation (oral)",
    "Conference presentation (poster)",
    "Web article",
]

def extract_report_date(filename: str) -> int:
    """Return the report date (YYYYMM) extracted from *filename*, or empty string if not found."""
    _REPORT_DATE_RE = re.compile(r"\d{6}")
    match = _REPORT_DATE_RE.search(filename)
    return int(match.group(0)) if match else None

def format_report_date(report_date: str) -> str:
    """Format a report date to YYYY-MM-DD format.
    Handles multiple input formats:
    - YYYYMM (6 digits): "202604" → "2026-04-01"
    - MM/DD/YYYY: "04/04/2026" → "2026-04-04"
    - Month DD, YYYY: "April 26, 2026" → "2026-04-26"
    """
    if not report_date:
        return report_date
    
    report_date = report_date.strip()
    # Try YYYYMM format (6 digits)
    if len(report_date) == 6 and report_date.isdigit():
        try:
            dt = datetime.strptime(f"{report_date}15", "%Y%m%d")
            return dt.strftime("%Y-%m-%d")
        except ValueError:
            pass
    # Try MM/DD/YYYY format
    try:
        dt = datetime.strptime(report_date, "%m/%d/%Y")
        return dt.strftime("%Y-%m-%d")
    except ValueError:
        pass
    # Try "Month DD, YYYY" format (e.g., "April 26, 2026")
    try:
        dt = datetime.strptime(report_date, "%B %d, %Y")
        return dt.strftime("%Y-%m-%d")
    except ValueError:
        pass
    # If none match, return as-is
    return report_date

def extract_project_id(filename: str) -> str:
    """Return the first project-ID token found in *filename*, or empty string."""
    _PROJECT_ID_RE = re.compile(r"\d{2}-[A-Z]+\d{2}-\d{4}")
    match = _PROJECT_ID_RE.search(filename)
    return match.group(0) if match else ""


def parse_table(table) -> list:
    """Return a list of rows (each a list of strings), skipping empty rows and de-duping merged cells."""
    parsed = []
    for row in table.rows:
        values = [cell.text.strip() for cell in row.cells]
        if not any(values):
            continue
        # deduped = [v for i, v in enumerate(values) if i == 0 or v != values[i - 1]]
        parsed.append(values)
    return parsed


def format_dict_table(data, title: str = "") -> str:
    """Render a flat or one-level-nested dict, or a category/sub_category/value DataFrame, as an ASCII table."""
    if isinstance(data, pd.DataFrame):
        rows = []
        last_cat = None
        for _, r in data.iterrows():
            cat = str(r.iloc[0])
            label = cat if cat != last_cat else ""
            last_cat = cat
            rows.append((label, str(r.iloc[1]), str(r.iloc[2])))
    else:
        rows = []
        for key, value in data.items():
            if isinstance(value, dict):
                first = True
                for sub_key, sub_val in value.items():
                    label = key if first else ""
                    rows.append((label, f"{sub_key}", str(sub_val)))
                    first = False
            else:
                rows.append((key, "", str(value)))

    col0 = max((len(r[0]) for r in rows), default=0)
    col1 = max((len(r[1]) for r in rows), default=0)
    col2 = max((len(r[2]) for r in rows), default=0)

    if col1 == 0:
        header = f"{'Category':<{col0}}  {'Value':<{col2}}"
        sep = f"{'-' * col0}  {'-' * col2}"
        lines = [f"{r[0]:<{col0}}  {r[2]:<{col2}}" for r in rows]
    else:
        header = f"{'Category':<{col0}}  {'Sub-category':<{col1}}  {'Value':<{col2}}"
        sep = f"{'-' * col0}  {'-' * col1}  {'-' * col2}"
        lines = [f"{r[0]:<{col0}}  {r[1]:<{col1}}  {r[2]:<{col2}}" for r in rows]

    parts = []
    if title:
        parts.append(title)
    parts += [header, sep] + lines
    return "\n".join(parts)


def check_volume_calculation(x: float, y: float, z: float, volume: float) -> bool:
    """Check that payload size (cm) and payload volume (U) are consistent.

    x * y * z in cm^3, divided by 1000, should equal volume in U (1U = 1000 cm^3).
    Returns True if consistent, False otherwise, and logs a warning on mismatch.
    """
    computed = x * y * z / 1000
    consistent = math.isclose(computed, volume)
    if not consistent:
        logger.warning(
            "Payload size %.4f x %.4f x %.4f cm gives %.4f U but payload volume is %.4f U",
            x,
            y,
            z,
            computed,
            volume,
        )
    return consistent


@dataclass(repr=False)
class Report:
    file_path: str | Path

    filename: str = field(init=False)
    project_id: str | None = field(init=False, default=None)
    summary: object = field(init=False, default=None)
    project_status: pd.DataFrame | None = field(init=False, default=None)
    summary: dict | None = field(init=False, default=None)
    slide_titles: list = field(init=False, default_factory=list)
    principal_investigator: str | None = field(init=False, default=None)
    accomodation_table: pd.DataFrame | None = field(init=False, default=None)
    publication_table: pd.DataFrame | None = field(init=False, default=None)
    patents_table: pd.DataFrame | None = field(init=False, default=None)
    student_metrics_table: pd.DataFrame | None = field(init=False, default=None)

    def __post_init__(self):
        # Ensure file_path is a Path object and check existence
        if isinstance(self.file_path, str):
            self.file_path = Path(self.file_path)
        self.filename = self.file_path.name
        if not self.file_path.exists():
            raise FileNotFoundError(f"Error: file not found: {self.file_path}")

        prs = Presentation(str(self.file_path))

        # iterate through the slides
        for idx, this_slide in enumerate(prs.slides, start=0):
            this_slide_title = self.get_slide_title(this_slide).lower()
            self.slide_titles.append(this_slide_title)
            if this_slide_title.lower() == "summary":
                self.summary = self.parse_summary_slide(this_slide)
                if isinstance(self.summary, dict):
                    self._process_summary(self.summary)
            if "Payload Accommodation".lower() in this_slide_title:
                self.accomodation_table = self.parse_accomodation_slide(this_slide)
            if "presentations" in this_slide_title:
                self.publication_table = self.parse_publications_slide(this_slide)
            if "patents" in this_slide_title:
                self.patents_table = self.parse_patents_slide(this_slide)
            if "student metrics and support" in this_slide_title:
                self.student_metrics_table = self.parse_student_metrics_slide(
                    this_slide
                )
            if "Project Summary".lower() in this_slide_title:
                self.project_status = self.parse_status_slide(this_slide)

    def parse_student_metrics_slide(self, slide) -> pd.DataFrame | None:
        """Extract the student metrics table information."""
        for shape in slide.shapes:
            if not shape.has_table:
                continue

            parsed_rows = parse_table(shape.table)
            if not parsed_rows:
                continue

            header = [h.strip() for h in parsed_rows[0]]
            data = parsed_rows[1:]
            if not data:
                return None
            return pd.DataFrame(data, columns=header)
        return None

    def _process_summary(self, summary: dict) -> None:
        """Extract fields from the parsed summary dict and validate values."""
        self.project_id = summary.get("proposal id", None)
        self.principal_investigator = summary.get("principal investigator", None)
        regime_value = summary.get("research regime", "") or ""
        valid = [r for r in RESEARCH_REGIMES if r in regime_value.lower()]
        if not valid:
            logger.warning(
                "%s: 'research regime' has no recognised value "
                "(got %r). Expected one or more of: %s",
                self.filename,
                regime_value,
                RESEARCH_REGIMES,
            )

    def parse_patents_slide(self, slide) -> pd.DataFrame | None:
        """Extract the patents table information."""
        for shape in slide.shapes:
            if not shape.has_table:
                continue

            parsed_rows = parse_table(shape.table)
            if not parsed_rows:
                continue

            header = [h.strip() for h in parsed_rows[0]]
            data = parsed_rows[1:]
            if not data:
                return None
            df = pd.DataFrame(data, columns=header)
            return df
        return None

    def parse_summary_slide(self, slide) -> dict:
        """Extract data from the project info slide"""
        for shape in slide.shapes:
            if not shape.has_table:
                continue

            parsed_rows = parse_table(shape.table)
            if not parsed_rows:
                continue
            result = {}
            for row in parsed_rows:
                if len(row) >= 2:
                    key = row[0].strip(":").lower()
                    value = row[1]
                    result[key] = value
            return result

        return None

    def parse_status_slide(self, slide) -> pd.DataFrame | None:
        """Extract data from the project status slide.

        Returns a DataFrame with columns [category, prior, current].
        """
        STATUS_CATEGORIES = {
            "cost",
            "schedule",
            "technical",
            "staffing",
            "facilities",
            "overall",
        }
        rows = []
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            parsed_rows = parse_table(shape.table)
            if not parsed_rows:
                continue
            for row in parsed_rows:
                if not row:
                    continue
                category = row[0].lower()
                if category in STATUS_CATEGORIES and len(row) >= 3:
                    prior = row[1].lower()
                    current = row[2].lower()
                    for label, value in (("prior", prior), ("current", current)):
                        if value not in STATUS_VALUES:
                            logger.warning(
                                "%s: status category %r has invalid %s value %r. "
                                "Expected one of: %s",
                                self.filename,
                                category,
                                label,
                                value,
                                STATUS_VALUES,
                            )
                    rows.append(
                        {"category": category, "prior": prior, "current": current}
                    )
        return (
            pd.DataFrame(rows, columns=["category", "prior", "current"])
            if rows
            else None
        )

    def parse_publications_slide(self, slide) -> pd.DataFrame | None:
        """Extract the publications table information."""
        for shape in slide.shapes:
            if not shape.has_table:
                continue

            parsed_rows = parse_table(shape.table)
            if not parsed_rows:
                continue

            header = [h.strip() for h in parsed_rows[0]]
            data = parsed_rows[1:]
            if not data:
                return None
            df = pd.DataFrame(data, columns=header)
            # Find the type column (first column whose header contains "type")
            type_col = next((c for c in df.columns if "type" in c.lower()), None)
            if type_col is not None:
                for val in df[type_col]:
                    if val and val not in PUBLICATION_TYPES:
                        logger.warning(
                            "%s: unrecognised publication type %r. Expected one of: %s",
                            self.filename,
                            val,
                            PUBLICATION_TYPES,
                        )
            return df
        return None

    def parse_accomodation_slide(self, slide) -> dict:
        """Extract the accomodation table information."""
        for shape in slide.shapes:
            if not shape.has_table:
                continue

            parsed_rows = parse_table(shape.table)
            if not parsed_rows:
                continue

            header = parsed_rows[0]
            try:
                cat_idx = next(
                    i
                    for i, h in enumerate(header)
                    if "requirement" in h.lower() and "category" in h.lower()
                )
                spec_idx = next(
                    (i for i, h in enumerate(header) if "specification" in h.lower()),
                    None,
                )
                unit_idx = next(
                    (i for i, h in enumerate(header) if "unit" in h.lower()), None
                )
                val_idx = next(i for i, h in enumerate(header) if "value" in h.lower())
            except StopIteration:
                logger.warning(
                    "Expected columns not found on accommodation slide. Skipping."
                )
                continue

            accom_rows = []
            for row in parsed_rows[1:]:
                if len(row) <= max(cat_idx, val_idx):
                    continue
                category = row[cat_idx].lower()
                raw_value = row[val_idx] if len(row) > val_idx else ""
                spec = (
                    row[spec_idx].strip()
                    if spec_idx is not None and len(row) > spec_idx
                    else ""
                )
                units = (
                    row[unit_idx].strip()
                    if unit_idx is not None and len(row) > unit_idx
                    else ""
                )

                value_parts = [v.strip() for v in raw_value.split(",")]

                if len(value_parts) > 1:
                    # Try to find labels from spec details, then units
                    spec_parts = [
                        s.strip().lower().split()[0]
                        for s in spec.split(",")
                        if s.strip()
                    ]
                    unit_parts = [
                        u.strip().lower().split()[0]
                        for u in units.split(",")
                        if u.strip()
                    ]
                    if len(spec_parts) == len(value_parts):
                        labels = spec_parts
                    elif len(unit_parts) == len(value_parts):
                        labels = unit_parts
                    else:
                        labels = None

                    if labels is not None:
                        for label, v in zip(labels, value_parts):
                            try:
                                accom_rows.append(
                                    {
                                        "category": category,
                                        "sub_category": label,
                                        "value": float(v),
                                    }
                                )
                            except ValueError:
                                accom_rows.append(
                                    {
                                        "category": category,
                                        "sub_category": label,
                                        "value": v.lower(),
                                    }
                                )
                    else:
                        accom_rows.append(
                            {
                                "category": category,
                                "sub_category": "",
                                "value": raw_value.lower(),
                            }
                        )
                else:
                    try:
                        accom_rows.append(
                            {
                                "category": category,
                                "sub_category": "",
                                "value": float(raw_value),
                            }
                        )
                    except ValueError:
                        accom_rows.append(
                            {
                                "category": category,
                                "sub_category": "",
                                "value": raw_value.lower(),
                            }
                        )

            return pd.DataFrame(
                accom_rows, columns=["category", "sub_category", "value"]
            )
        return None

    def get_slide_title(self, slide) -> str:
        """Return the title text of a slide, or an empty string if none found."""
        if slide.shapes.title is not None:
            return slide.shapes.title.text or ""
        # Fall back to checking all placeholder shapes for a title placeholder
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 0:
                return shape.text or ""
        return ""

    def __repr__(self):
        _STATUS_COLOR = {
            "g": "\033[32m",
            "green": "\033[32m",
            "y": "\033[33m",
            "yellow": "\033[33m",
            "r": "\033[31m",
            "red": "\033[31m",
        }
        _RESET = "\033[0m"

        def colorize(value: str) -> str:
            color = _STATUS_COLOR.get(value, "")
            return f"{color}{value.upper()}{_RESET}" if color else value.upper()

        result = f"Report {self.filename} ({len(self.slide_titles)} slides)\n"
        result += f"Project ID: {self.project_id}, PI: {self.principal_investigator}\n"
        if self.project_status is not None and not self.project_status.empty:
            df = self.project_status
            # Put overall first, then the rest in natural order
            mask = df["category"] == "overall"
            df = pd.concat([df[mask], df[~mask]], ignore_index=True)

            col0 = df["category"].str.len().max()
            col1 = 7  # visible width for Prior column
            hdr = f"\n{'Category':<{col0}}  {'Prior':<{col1}}  {'Current'}"
            sep = f"{'-' * col0}  {'-' * col1}  {'-' * col1}"
            lines = [hdr, sep]
            for _, row in df.iterrows():
                prior = colorize(row["prior"])
                current = colorize(row["current"])
                prior_pad = prior + " " * max(0, col1 - len(row["prior"].upper()))
                lines.append(f"{row['category']:<{col0}}  {prior_pad}  {current}")
            result += "\n".join(lines) + "\n"
        result += "\n"
        if self.accomodation_table is not None and not self.accomodation_table.empty:
            result += (
                format_dict_table(self.accomodation_table, title="Accommodation:")
                + "\n"
            )
        if self.publication_table is not None and not self.publication_table.empty:
            result += (
                f"\nPublications:\n{self.publication_table.to_string(index=False)}\n"
            )
        return result


def get_accomodation_data(reports: list[Report]) -> None:
    # Collect all accommodation keys across all reports
    all_keys = []
    seen_keys = set()
    for report in reports:
        if report.accomodation_table is not None:
            for _, r in report.accomodation_table.iterrows():
                k = (
                    f"{r['category']} {r['sub_category']}".strip()
                    if r["sub_category"]
                    else r["category"]
                )
                if k not in seen_keys:
                    all_keys.append(k)
                    seen_keys.add(k)

    # Write one row per project to CSV
    csv_path = Path("payload_accommodation.csv")
    with csv_path.open("w", newline="", encoding="utf-8") as fh:
        fieldnames = ["filename", "project_id", "principal_investigator"] + all_keys
        writer = csv.DictWriter(fh, fieldnames=fieldnames)
        writer.writeheader()
        for report in reports:
            row = {
                "filename": report.filename,
                "project_id": report.project_id or "",
                "principal_investigator": report.principal_investigator or "",
            }
            if report.accomodation_table is not None:
                for _, r in report.accomodation_table.iterrows():
                    k = (
                        f"{r['category']} {r['sub_category']}".strip()
                        if r["sub_category"]
                        else r["category"]
                    )
                    row[k] = r["value"]
            writer.writerow(row)
    print(f"Wrote accommodation table to {csv_path}")
